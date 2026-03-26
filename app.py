import os
import csv
import uuid
from io import StringIO, BytesIO
from datetime import datetime, date, timedelta
from functools import wraps
from flask import Flask, render_template, url_for, flash, redirect, request, send_file, abort, jsonify
from flask_login import LoginManager, login_user, current_user, logout_user, login_required
from werkzeug.utils import secure_filename
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from pptx import Presentation
from flask_mail import Mail, Message
from sqlalchemy import func, extract

# Import your models and forms
from models import db, User, Driver, Task, Company, Vehicle, Payment, Invoice, MaintenanceRecord
from forms import (
    RegistrationForm, LoginForm, DriverForm, TaskForm, SearchForm,
    VehicleForm, PaymentForm, InvoiceForm, AdminUserEditForm
)

# Helper for route calculations (optional)
# from route_utils import get_distance_and_time  # uncomment if you have API key

app = Flask(__name__)

# ----------------------------------------------------------------------
# Configuration (using environment variables, no separate config.py)
# ----------------------------------------------------------------------
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'default-secret-key')
app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL',
    'sqlite:///' + os.path.join(os.path.abspath(os.path.dirname(__file__)), 'instance', 'logistics.db'))
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

# Upload folder for driver documents
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.abspath(os.path.dirname(__file__)), 'static', 'uploads')
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Mail settings (for notifications)
app.config['MAIL_SERVER'] = 'smtp.gmail.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USERNAME'] = os.environ.get('MAIL_USERNAME', 'levido14@gmail.com')
app.config['MAIL_PASSWORD'] = os.environ.get('MAIL_PASSWORD', '')

# OpenRouteService API key (optional)
app.config['ORS_API_KEY'] = os.environ.get('ORS_API_KEY', '')

# Twilio settings (optional)
app.config['TWILIO_ACCOUNT_SID'] = os.environ.get('TWILIO_ACCOUNT_SID', '')
app.config['TWILIO_AUTH_TOKEN'] = os.environ.get('TWILIO_AUTH_TOKEN', '')
app.config['TWILIO_PHONE'] = os.environ.get('TWILIO_PHONE', '')

# Initialize extensions
db.init_app(app)
mail = Mail(app)
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

# ----------------------------------------------------------------------
# Database initialization (runs on app start, e.g., when gunicorn loads)
# ----------------------------------------------------------------------
with app.app_context():
    db.create_all()
    if Company.query.count() == 0:
        default_company = Company(name='Default Company')
        db.session.add(default_company)
        db.session.commit()
        print("Default company created")
    else:
        print("Company already exists")

# ----------------------------------------------------------------------
# User loader
# ----------------------------------------------------------------------
@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# ----------------------------------------------------------------------
# Admin decorator
# ----------------------------------------------------------------------
def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role != 'admin':
            abort(403)
        return f(*args, **kwargs)
    return decorated_function

# ----------------------------------------------------------------------
# File handling
# ----------------------------------------------------------------------
def save_file(file, prefix=''):
    if file and file.filename:
        ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else ''
        unique_name = f"{prefix}_{uuid.uuid4().hex}.{ext}" if prefix else f"{uuid.uuid4().hex}.{ext}"
        folder = app.config['UPLOAD_FOLDER']
        os.makedirs(folder, exist_ok=True)
        filepath = os.path.join(folder, unique_name)
        try:
            file.save(filepath)
            return os.path.join('uploads', unique_name).replace('\\', '/')
        except Exception as e:
            print(f"Error saving file: {e}")
            return None
    return None

# ----------------------------------------------------------------------
# Email and SMS helpers (uncomment when ready)
# ----------------------------------------------------------------------
def send_email_task_assigned(task, assigner):
    subject = f"New Task Assigned to {task.driver.name}"
    body = f"""
    Task ID: {task.id}
    Driver: {task.driver.name} (Car: {task.driver.car_number_plate})
    Load Location: {task.load_location}
    Offload Location: {task.offload_location}
    Assigned by: {assigner.username}
    """
    recipients = [assigner.email]
    # Uncomment to send email
    # msg = Message(subject, recipients=recipients)
    # msg.body = body
    # mail.send(msg)
    print(f"Email would be sent to {recipients}:\n{body}")

def send_email_task_completed(task, completer):
    subject = f"Task Completed: {task.driver.name}"
    body = f"""
    Task ID: {task.id}
    Driver: {task.driver.name} (Car: {task.driver.car_number_plate})
    Load Location: {task.load_location}
    Offload Location: {task.offload_location}
    Completed by: {completer.username}
    Completed on: {task.completed_date}
    """
    recipients = [completer.email]
    # Uncomment to send email
    # msg = Message(subject, recipients=recipients)
    # msg.body = body
    # mail.send(msg)
    print(f"Email would be sent to {recipients}:\n{body}")

def send_sms(to, body):
    if app.config['TWILIO_ACCOUNT_SID'] and app.config['TWILIO_AUTH_TOKEN']:
        from twilio.rest import Client
        client = Client(app.config['TWILIO_ACCOUNT_SID'], app.config['TWILIO_AUTH_TOKEN'])
        client.messages.create(body=body, from_=app.config['TWILIO_PHONE'], to=to)
    else:
        print(f"SMS would be sent to {to}: {body}")

# ----------------------------------------------------------------------
# Report generation functions (unchanged)
# ----------------------------------------------------------------------
def generate_driver_csv(driver, tasks):
    output = StringIO()
    writer = csv.writer(output)
    writer.writerow(['Driver Name', driver.name])
    writer.writerow(['Car Plate', driver.car_number_plate])
    writer.writerow(['Phone', driver.phone])
    writer.writerow(['Tasks:'])
    writer.writerow(['Date', 'Load Location', 'Offload Location', 'Completed'])
    for task in tasks:
        writer.writerow([task.assigned_date, task.load_location, task.offload_location, 'Yes' if task.completed else 'No'])
    output.seek(0)
    return send_file(BytesIO(output.getvalue().encode('utf-8')), mimetype='text/csv', as_attachment=True, download_name=f'{driver.car_number_plate}_report.csv')

def generate_driver_pdf(driver, tasks):
    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(f"Driver Report: {driver.name}", styles['Title']))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"Car Plate: {driver.car_number_plate}", styles['Normal']))
    story.append(Paragraph(f"Phone: {driver.phone}", styles['Normal']))
    story.append(Spacer(1, 12))
    data = [['Date', 'Load Location', 'Offload Location', 'Completed']]
    for task in tasks:
        data.append([str(task.assigned_date), task.load_location, task.offload_location, 'Yes' if task.completed else 'No'])
    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.grey),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 12),
        ('BACKGROUND', (0,1), (-1,-1), colors.beige),
        ('GRID', (0,0), (-1,-1), 1, colors.black)
    ]))
    story.append(table)
    doc.build(story)
    output.seek(0)
    return send_file(output, mimetype='application/pdf', as_attachment=True, download_name=f'{driver.car_number_plate}_report.pdf')

def generate_driver_ppt(driver, tasks):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = f"Driver Report: {driver.name}"
    text = f"Car Plate: {driver.car_number_plate}\nPhone: {driver.phone}\n\nTasks:\n"
    for task in tasks:
        status = 'Completed' if task.completed else 'Pending'
        text += f"- {task.assigned_date}: {task.load_location} → {task.offload_location} ({status})\n"
    content.text = text
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name=f'{driver.car_number_plate}_report.pptx')

def generate_summary_csv(drivers):
    output = StringIO()
    writer = csv.writer(output)
    writer.writerow(['Driver Name', 'Car Plate', 'Total Tasks', 'Completed', 'Pending'])
    for driver in drivers:
        tasks = Task.query.filter_by(driver_id=driver.id).all()
        total = len(tasks)
        completed = sum(1 for t in tasks if t.completed)
        pending = total - completed
        writer.writerow([driver.name, driver.car_number_plate, total, completed, pending])
    output.seek(0)
    return send_file(BytesIO(output.getvalue().encode('utf-8')), mimetype='text/csv', as_attachment=True, download_name='summary_report.csv')

def generate_summary_pdf(drivers):
    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph("Overall Summary Report", styles['Title']))
    story.append(Spacer(1, 12))
    data = [['Driver Name', 'Car Plate', 'Total Tasks', 'Completed', 'Pending']]
    for driver in drivers:
        tasks = Task.query.filter_by(driver_id=driver.id).all()
        total = len(tasks)
        completed = sum(1 for t in tasks if t.completed)
        pending = total - completed
        data.append([driver.name, driver.car_number_plate, str(total), str(completed), str(pending)])
    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.grey),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 12),
        ('BACKGROUND', (0,1), (-1,-1), colors.beige),
        ('GRID', (0,0), (-1,-1), 1, colors.black)
    ]))
    story.append(table)
    doc.build(story)
    output.seek(0)
    return send_file(output, mimetype='application/pdf', as_attachment=True, download_name='summary_report.pdf')

def generate_summary_ppt(drivers):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Overall Summary Report"
    text = ""
    for driver in drivers:
        tasks = Task.query.filter_by(driver_id=driver.id).all()
        total = len(tasks)
        completed = sum(1 for t in tasks if t.completed)
        pending = total - completed
        text += f"{driver.name} ({driver.car_number_plate}): {completed}/{total} completed, {pending} pending\n"
    content.text = text
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name='summary_report.pptx')

# ----------------------------------------------------------------------
# Routes
# ----------------------------------------------------------------------
@app.route('/')
def index():
    return render_template('index.html')

# Registration (keep but you may want to restrict later)
@app.route('/register', methods=['GET', 'POST'])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    form = RegistrationForm()
    user_count = User.query.count()
    if form.validate_on_submit():
        # In production, admin would assign company. For simplicity, use company_id=1.
        role = form.role.data
        user = User(username=form.username.data, email=form.email.data, role=role, company_id=1)
        user.set_password(form.password.data)
        db.session.add(user)
        db.session.commit()
        flash('Your account has been created! You can now log in.', 'success')
        return redirect(url_for('login'))
    return render_template('register.html', form=form, user_count=user_count)

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(username=form.username.data).first()
        if user and user.check_password(form.password.data):
            login_user(user, remember=form.remember.data)
            next_page = request.args.get('next')
            return redirect(next_page) if next_page else redirect(url_for('dashboard'))
        else:
            flash('Login unsuccessful. Check username and password.', 'danger')
    return render_template('login.html', form=form)

@app.route('/logout')
def logout():
    logout_user()
    return redirect(url_for('index'))

# Dashboard
@app.route('/dashboard')
@login_required
def dashboard():
    if current_user.role == 'admin':
        drivers = Driver.query.filter_by(active=True, company_id=current_user.company_id).all()
        tasks = Task.query.filter_by(company_id=current_user.company_id).all()
        users = User.query.filter_by(company_id=current_user.company_id).all()
    else:
        drivers = Driver.query.filter_by(user_id=current_user.id, active=True).all()
        tasks = Task.query.filter_by(user_id=current_user.id).all()
        users = None
    return render_template('dashboard.html', drivers=drivers, tasks=tasks, users=users)

# API endpoint for dashboard analytics
@app.route('/api/dashboard_data')
@login_required
def dashboard_data():
    tasks = Task.query.filter_by(company_id=current_user.company_id)
    drivers = Driver.query.filter_by(company_id=current_user.company_id, active=True)
    tasks_completed = tasks.filter_by(completed=True).count()
    tasks_pending = tasks.filter_by(completed=False).count()
    drivers_count = drivers.count()
    monthly = db.session.query(
        extract('year', Task.completed_date).label('year'),
        extract('month', Task.completed_date).label('month'),
        func.count(Task.id).label('count')
    ).filter(Task.completed == True, Task.company_id == current_user.company_id) \
     .group_by('year', 'month') \
     .order_by('year', 'month') \
     .limit(6).all()
    monthly_data = [{'month': f"{int(m[1])}/{int(m[0])}", 'count': m[2]} for m in monthly]
    return jsonify({
        'tasks_completed': tasks_completed,
        'tasks_pending': tasks_pending,
        'drivers_count': drivers_count,
        'monthly_data': monthly_data
    })

# Driver management
@app.route('/driver/add', methods=['GET', 'POST'])
@login_required
def add_driver():
    form = DriverForm()
    # Populate vehicle choices for current company
    form.vehicle_id.choices = [(0, 'None')] + [(v.id, f"{v.make} {v.model} ({v.plate})") for v in Vehicle.query.filter_by(company_id=current_user.company_id).all()]
    if form.validate_on_submit():
        cert_path = save_file(form.certificate_of_conduct.data, prefix='cert')
        license_path = save_file(form.driving_license.data, prefix='license')
        id_path = save_file(form.national_id.data, prefix='id')
        vehicle_path = save_file(form.vehicle_picture.data, prefix='vehicle')
        passport_path = save_file(form.passport_photo.data, prefix='passport')
        driver = Driver(
            name=form.name.data,
            phone=form.phone.data,
            car_number_plate=form.car_number_plate.data,
            certificate_of_conduct=cert_path,
            driving_license=license_path,
            national_id=id_path,
            vehicle_picture=vehicle_path,
            passport_photo=passport_path,
            user_id=current_user.id,
            company_id=current_user.company_id,
            vehicle_id=form.vehicle_id.data if form.vehicle_id.data != 0 else None
        )
        db.session.add(driver)
        db.session.commit()
        flash('Driver registered successfully!', 'success')
        return redirect(url_for('dashboard'))
    return render_template('add_driver.html', form=form)

@app.route('/driver/<int:id>')
@login_required
def driver_detail(id):
    driver = Driver.query.get_or_404(id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    return render_template('driver_detail.html', driver=driver)

@app.route('/driver/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit_driver(id):
    driver = Driver.query.get_or_404(id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    form = DriverForm()
    # Populate vehicle choices
    form.vehicle_id.choices = [(0, 'None')] + [(v.id, f"{v.make} {v.model} ({v.plate})") for v in Vehicle.query.filter_by(company_id=current_user.company_id).all()]
    if form.validate_on_submit():
        driver.name = form.name.data
        driver.phone = form.phone.data
        driver.car_number_plate = form.car_number_plate.data
        driver.vehicle_id = form.vehicle_id.data if form.vehicle_id.data != 0 else None
        if form.certificate_of_conduct.data:
            driver.certificate_of_conduct = save_file(form.certificate_of_conduct.data, prefix='cert')
        if form.driving_license.data:
            driver.driving_license = save_file(form.driving_license.data, prefix='license')
        if form.national_id.data:
            driver.national_id = save_file(form.national_id.data, prefix='id')
        if form.vehicle_picture.data:
            driver.vehicle_picture = save_file(form.vehicle_picture.data, prefix='vehicle')
        if form.passport_photo.data:
            driver.passport_photo = save_file(form.passport_photo.data, prefix='passport')
        db.session.commit()
        flash('Driver details updated.', 'success')
        return redirect(url_for('driver_detail', id=driver.id))
    # Pre-populate
    form.name.data = driver.name
    form.phone.data = driver.phone
    form.car_number_plate.data = driver.car_number_plate
    form.vehicle_id.data = driver.vehicle_id or 0
    return render_template('edit_driver.html', form=form, driver=driver)

@app.route('/driver/<int:id>/sack')
@login_required
def sack_driver(id):
    driver = Driver.query.get_or_404(id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    driver.active = False
    db.session.commit()
    flash(f'Driver {driver.name} has been sacked.', 'warning')
    return redirect(url_for('dashboard'))

@app.route('/driver/<int:id>/delete_doc/<doc_type>', methods=['POST'])
@login_required
def delete_doc(id, doc_type):
    driver = Driver.query.get_or_404(id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    fields = {
        'cert': 'certificate_of_conduct',
        'license': 'driving_license',
        'id': 'national_id',
        'vehicle': 'vehicle_picture',
        'passport': 'passport_photo'
    }
    field = fields.get(doc_type)
    if field:
        file_path = getattr(driver, field)
        if file_path:
            full_path = os.path.join(app.config['UPLOAD_FOLDER'], os.path.basename(file_path))
            if os.path.exists(full_path):
                os.remove(full_path)
            setattr(driver, field, None)
            db.session.commit()
            flash('Document deleted.', 'success')
    return redirect(url_for('driver_detail', id=driver.id))

# Vehicle management
@app.route('/vehicle/add', methods=['GET', 'POST'])
@login_required
def add_vehicle():
    form = VehicleForm()
    if form.validate_on_submit():
        vehicle = Vehicle(
            make=form.make.data,
            model=form.model.data,
            plate=form.plate.data,
            capacity_kg=form.capacity_kg.data,
            insurance_expiry=form.insurance_expiry.data,
            status=form.status.data,
            company_id=current_user.company_id
        )
        db.session.add(vehicle)
        db.session.commit()
        flash('Vehicle added successfully!', 'success')
        return redirect(url_for('list_vehicles'))
    return render_template('add_vehicle.html', form=form)

@app.route('/vehicles')
@login_required
def list_vehicles():
    vehicles = Vehicle.query.filter_by(company_id=current_user.company_id).all()
    return render_template('list_vehicles.html', vehicles=vehicles)

@app.route('/vehicle/<int:id>')
@login_required
def vehicle_detail(id):
    vehicle = Vehicle.query.get_or_404(id)
    if vehicle.company_id != current_user.company_id:
        abort(403)
    return render_template('vehicle_detail.html', vehicle=vehicle)

@app.route('/vehicle/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit_vehicle(id):
    vehicle = Vehicle.query.get_or_404(id)
    if vehicle.company_id != current_user.company_id:
        abort(403)
    form = VehicleForm()
    if form.validate_on_submit():
        vehicle.make = form.make.data
        vehicle.model = form.model.data
        vehicle.plate = form.plate.data
        vehicle.capacity_kg = form.capacity_kg.data
        vehicle.insurance_expiry = form.insurance_expiry.data
        vehicle.status = form.status.data
        db.session.commit()
        flash('Vehicle updated.', 'success')
        return redirect(url_for('vehicle_detail', id=vehicle.id))
    # Pre-populate
    form.make.data = vehicle.make
    form.model.data = vehicle.model
    form.plate.data = vehicle.plate
    form.capacity_kg.data = vehicle.capacity_kg
    form.insurance_expiry.data = vehicle.insurance_expiry
    form.status.data = vehicle.status
    return render_template('edit_vehicle.html', form=form, vehicle=vehicle)

# Task management (with distance calculation placeholder)
@app.route('/driver/<int:driver_id>/task/add', methods=['GET', 'POST'])
@login_required
def add_task(driver_id):
    driver = Driver.query.get_or_404(driver_id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    form = TaskForm()
    if form.validate_on_submit():
        # Optionally calculate distance (uncomment if API key set)
        # distance_km, duration_sec = None, None
        # try:
        #     from route_utils import get_distance_and_time
        #     distance_km, duration_sec = get_distance_and_time(form.load_location.data, form.offload_location.data)
        # except:
        #     pass
        # fuel_cost = distance_km * 0.3 if distance_km else None
        task = Task(
            load_location=form.load_location.data,
            offload_location=form.offload_location.data,
            driver_id=driver.id,
            user_id=current_user.id,
            company_id=current_user.company_id,
            # distance_km=distance_km,
            # estimated_duration_sec=duration_sec,
            # fuel_cost=fuel_cost
        )
        db.session.add(task)
        db.session.commit()
        send_email_task_assigned(task, current_user)
        if driver.phone:
            send_sms(driver.phone, f"New task assigned: {task.load_location} to {task.offload_location}")
        flash('Task assigned successfully!', 'success')
        return redirect(url_for('driver_detail', id=driver.id))
    return render_template('add_task.html', form=form, driver=driver)

@app.route('/task/<int:id>/complete')
@login_required
def complete_task(id):
    task = Task.query.get_or_404(id)
    driver = Driver.query.get(task.driver_id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    task.completed = True
    task.completed_date = datetime.utcnow()
    db.session.commit()
    send_email_task_completed(task, current_user)
    if driver.phone:
        send_sms(driver.phone, f"Task completed: {task.load_location} to {task.offload_location}")
    flash('Task marked as completed.', 'success')
    return redirect(url_for('driver_detail', id=task.driver_id))

# Advanced search
@app.route('/search', methods=['GET', 'POST'])
@login_required
def search():
    form = SearchForm()
    query = Driver.query.filter_by(active=True, company_id=current_user.company_id)
    if request.args.get('plate'):
        query = query.filter(Driver.car_number_plate.ilike(f"%{request.args['plate']}%"))
    if request.args.get('name'):
        query = query.filter(Driver.name.ilike(f"%{request.args['name']}%"))
    if request.args.get('vehicle_status'):
        query = query.join(Vehicle).filter(Vehicle.status == request.args['vehicle_status'])
    if request.args.get('expiry_soon'):
        soon = date.today() + timedelta(days=30)
        query = query.join(Vehicle).filter(Vehicle.insurance_expiry <= soon)
    drivers = query.all()
    return render_template('search_results.html', form=form, drivers=drivers)

# Payment management
@app.route('/driver/<int:driver_id>/payment/add', methods=['GET', 'POST'])
@login_required
def add_payment(driver_id):
    driver = Driver.query.get_or_404(driver_id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    form = PaymentForm()
    form.task_id.choices = [(0, 'None')] + [(t.id, f"{t.load_location} → {t.offload_location}") for t in Task.query.filter_by(driver_id=driver.id).all()]
    if form.validate_on_submit():
        payment = Payment(
            driver_id=driver.id,
            amount=form.amount.data,
            type=form.type.data,
            description=form.description.data,
            task_id=form.task_id.data if form.task_id.data != 0 else None
        )
        db.session.add(payment)
        db.session.commit()
        flash('Payment recorded.', 'success')
        return redirect(url_for('driver_detail', id=driver.id))
    return render_template('add_payment.html', form=form, driver=driver)

# Invoice management
@app.route('/invoice/create', methods=['GET', 'POST'])
@login_required
def create_invoice():
    form = InvoiceForm()
    tasks = Task.query.filter_by(company_id=current_user.company_id, completed=True).all()
    form.tasks.choices = [(t.id, f"{t.driver.name}: {t.load_location} → {t.offload_location}") for t in tasks]
    if form.validate_on_submit():
        selected_tasks = Task.query.filter(Task.id.in_(form.tasks.data)).all()
        total = sum(task.fuel_cost or 0 for task in selected_tasks)  # adjust as needed
        invoice = Invoice(
            client_name=form.client_name.data,
            client_email=form.client_email.data,
            total=total,
            paid=False,
            company_id=current_user.company_id,
            tasks=selected_tasks
        )
        db.session.add(invoice)
        db.session.commit()
        flash('Invoice created.', 'success')
        return redirect(url_for('list_invoices'))
    return render_template('create_invoice.html', form=form)

@app.route('/invoices')
@login_required
def list_invoices():
    invoices = Invoice.query.filter_by(company_id=current_user.company_id).all()
    return render_template('list_invoices.html', invoices=invoices)

@app.route('/invoice/<int:id>/pdf')
@login_required
def invoice_pdf(id):
    invoice = Invoice.query.get_or_404(id)
    if invoice.company_id != current_user.company_id:
        abort(403)
    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(f"Invoice #{invoice.id}", styles['Title']))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"Client: {invoice.client_name}", styles['Normal']))
    story.append(Paragraph(f"Date: {invoice.date.strftime('%Y-%m-%d')}", styles['Normal']))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Tasks:", styles['Heading2']))
    data = [['Driver', 'Load → Offload', 'Amount']]
    for task in invoice.tasks:
        data.append([task.driver.name, f"{task.load_location} → {task.offload_location}", f"${task.fuel_cost or 0:.2f}"])
    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.grey),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 12),
        ('BACKGROUND', (0,1), (-1,-1), colors.beige),
        ('GRID', (0,0), (-1,-1), 1, colors.black)
    ]))
    story.append(table)
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"Total: ${invoice.total:.2f}", styles['Normal']))
    doc.build(story)
    output.seek(0)
    return send_file(output, mimetype='application/pdf', as_attachment=True, download_name=f'invoice_{invoice.id}.pdf')

# Admin user management
@app.route('/admin/users')
@login_required
@admin_required
def list_users():
    users = User.query.filter_by(company_id=current_user.company_id).all()
    return render_template('admin/users.html', users=users)

@app.route('/admin/user/create', methods=['GET', 'POST'])
@login_required
@admin_required
def create_user():
    form = AdminUserEditForm()
    if form.validate_on_submit():
        user = User(
            username=form.username.data,
            email=form.email.data,
            role=form.role.data,
            company_id=current_user.company_id
        )
        user.set_password(form.password.data)
        db.session.add(user)
        db.session.commit()
        flash('User created.', 'success')
        return redirect(url_for('list_users'))
    return render_template('admin/create_user.html', form=form)

@app.route('/admin/user/<int:id>/edit', methods=['GET', 'POST'])
@login_required
@admin_required
def edit_user(id):
    user = User.query.get_or_404(id)
    if user.company_id != current_user.company_id:
        abort(403)
    form = AdminUserEditForm()
    if form.validate_on_submit():
        user.username = form.username.data
        user.email = form.email.data
        user.role = form.role.data
        if form.password.data:
            user.set_password(form.password.data)
        db.session.commit()
        flash('User updated.', 'success')
        return redirect(url_for('list_users'))
    form.username.data = user.username
    form.email.data = user.email
    form.role.data = user.role
    return render_template('admin/edit_user.html', form=form, user=user)

@app.route('/admin/user/<int:id>/delete', methods=['POST'])
@login_required
@admin_required
def delete_user(id):
    user = User.query.get_or_404(id)
    if user.id == current_user.id:
        flash('You cannot delete yourself.', 'danger')
        return redirect(url_for('list_users'))
    if user.company_id != current_user.company_id:
        abort(403)
    db.session.delete(user)
    db.session.commit()
    flash('User deleted.', 'success')
    return redirect(url_for('list_users'))

# Reports
@app.route('/report/driver/<int:driver_id>')
@login_required
def driver_report(driver_id):
    driver = Driver.query.get_or_404(driver_id)
    if current_user.role != 'admin' and driver.user_id != current_user.id:
        abort(403)
    fmt = request.args.get('format', 'csv')
    tasks = Task.query.filter_by(driver_id=driver.id).order_by(Task.assigned_date).all()
    if fmt == 'csv':
        return generate_driver_csv(driver, tasks)
    elif fmt == 'pdf':
        return generate_driver_pdf(driver, tasks)
    elif fmt == 'ppt':
        return generate_driver_ppt(driver, tasks)
    else:
        abort(400)

@app.route('/report/summary')
@login_required
def summary_report():
    fmt = request.args.get('format', 'csv')
    if current_user.role == 'admin':
        drivers = Driver.query.filter_by(active=True, company_id=current_user.company_id).all()
    else:
        drivers = Driver.query.filter_by(user_id=current_user.id, active=True).all()
    if fmt == 'csv':
        return generate_summary_csv(drivers)
    elif fmt == 'pdf':
        return generate_summary_pdf(drivers)
    elif fmt == 'ppt':
        return generate_summary_ppt(drivers)
    else:
        abort(400)

# ----------------------------------------------------------------------
# Temporary admin creation route (remove after first use)
# ----------------------------------------------------------------------
@app.route('/create-admin')
def create_admin():
    with app.app_context():
        # Ensure default company exists
        company = Company.query.first()
        if not company:
            company = Company(name='Default Company')
            db.session.add(company)
            db.session.commit()
        # Create admin user if not exists
        if not User.query.filter_by(username='admin').first():
            admin = User(username='admin', email='admin@example.com', role='admin', company_id=company.id)
            admin.set_password('admin')
            db.session.add(admin)
            db.session.commit()
            return "Admin user created. You can now login with admin/admin."
        else:
            return "Admin user already exists."

# ----------------------------------------------------------------------
# Main entry point (for local development)
# ----------------------------------------------------------------------
if __name__ == '__main__':
    app.run(debug=True)